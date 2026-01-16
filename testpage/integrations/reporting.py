from __future__ import annotations

from datetime import datetime
from io import BytesIO


def _escape_pdf_text(text: str) -> str:
    return text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")


def _build_minimal_pdf(lines: list[str]) -> bytes:
    y_position = 760
    line_height = 16
    text_lines = []
    for line in lines:
        safe = _escape_pdf_text(line)
        text_lines.append(f"1 0 0 1 50 {y_position} Tm ({safe}) Tj")
        y_position -= line_height
    content_stream = "BT /F1 12 Tf\n" + "\n".join(text_lines) + "\nET"
    content_bytes = content_stream.encode("latin1")

    objects = []
    objects.append(b"1 0 obj << /Type /Catalog /Pages 2 0 R >> endobj")
    objects.append(b"2 0 obj << /Type /Pages /Kids [3 0 R] /Count 1 >> endobj")
    objects.append(
        b"3 0 obj << /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >> endobj"
    )
    objects.append(
        f"4 0 obj << /Length {len(content_bytes)} >> stream\n".encode("latin1")
        + content_bytes
        + b"\nendstream endobj"
    )
    objects.append(b"5 0 obj << /Type /Font /Subtype /Type1 /BaseFont /Helvetica >> endobj")

    xref_positions = []
    output = BytesIO()
    output.write(b"%PDF-1.4\n")
    for obj in objects:
        xref_positions.append(output.tell())
        output.write(obj + b"\n")

    xref_start = output.tell()
    output.write(f"xref\n0 {len(objects) + 1}\n".encode("latin1"))
    output.write(b"0000000000 65535 f \n")
    for pos in xref_positions:
        output.write(f"{pos:010d} 00000 n \n".encode("latin1"))
    output.write(
        f"trailer << /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref_start}\n%%EOF".encode("latin1")
    )
    return output.getvalue()


def apply_text_mapping(presentation, mapping: dict[str, str]) -> None:
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for key, value in mapping.items():
                        if key in run.text:
                            run.text = run.text.replace(key, value)


def build_pptx_report(
    title: str,
    summary_lines: list[str],
    highlights: list[str],
    *,
    template_path: str | None = None,
) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation(template_path) if template_path else Presentation()

    if not template_path:
        presentation.slide_width = Inches(13.33)
        presentation.slide_height = Inches(7.5)

        title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        title_slide.shapes.title.text = "{{TITLE}}"
        title_slide.placeholders[1].text = "{{SUBTITLE}}"

        summary_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        summary_slide.shapes.title.text = "Executive Summary"
        summary_box = summary_slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.5), Inches(4.5))
        summary_box.text_frame.text = "{{SUMMARY}}"

        highlight_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        highlight_slide.shapes.title.text = "Highlights"
        highlight_box = highlight_slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.5), Inches(4.5))
        highlight_box.text_frame.text = "{{HIGHLIGHTS}}"

    mapping = {
        "{{TITLE}}": title,
        "{{SUBTITLE}}": "Click Insight Hub | HR Analytics",
        "{{SUMMARY}}": "\n".join(summary_lines) if summary_lines else "요약 정보가 없습니다.",
        "{{HIGHLIGHTS}}": "\n".join(highlights) if highlights else "하이라이트 정보가 없습니다.",
        "{{DATE}}": datetime.now().strftime("%Y-%m-%d"),
    }
    apply_text_mapping(presentation, mapping)

    output = BytesIO()
    presentation.save(output)
    output.seek(0)
    return output.getvalue()


def build_pdf_report(title: str, summary_lines: list[str], highlights: list[str]) -> bytes:
    lines = [
        title,
        "Click Insight Hub | HR Analytics",
        f"Report Date: {datetime.now().strftime('%Y-%m-%d')}",
        "",
        "Executive Summary",
        *summary_lines,
        "",
        "Highlights",
        *highlights,
    ]

    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas

        output = BytesIO()
        pdf = canvas.Canvas(output, pagesize=letter)
        textobject = pdf.beginText(50, 750)
        for line in lines:
            textobject.textLine(line)
        pdf.drawText(textobject)
        pdf.showPage()
        pdf.save()
        output.seek(0)
        return output.getvalue()
    except Exception:
        return _build_minimal_pdf(lines)


def build_excel_report(title: str, summary_lines: list[str], highlights: list[str]) -> bytes:
    try:
        import pandas as pd

        summary_df = pd.DataFrame({"Summary": summary_lines or ["요약 정보가 없습니다."]})
        highlights_df = pd.DataFrame({"Highlights": highlights or ["하이라이트 정보가 없습니다."]})
        meta_df = pd.DataFrame(
            {
                "Title": [title],
                "Generated At": [datetime.now().strftime("%Y-%m-%d %H:%M:%S")],
                "Source": ["Click Insight Hub | HR Analytics"],
            }
        )

        output = BytesIO()
        with pd.ExcelWriter(output, engine="openpyxl") as writer:
            meta_df.to_excel(writer, sheet_name="Meta", index=False)
            summary_df.to_excel(writer, sheet_name="Summary", index=False)
            highlights_df.to_excel(writer, sheet_name="Highlights", index=False)
        output.seek(0)
        return output.getvalue()
    except Exception:
        fallback = BytesIO()
        fallback.write("Title,Section,Text\n".encode("utf-8"))
        fallback.write(f"{title},Summary,".encode("utf-8"))
        fallback.write(" | ".join(summary_lines).encode("utf-8"))
        fallback.write("\n".encode("utf-8"))
        fallback.write(f"{title},Highlights,".encode("utf-8"))
        fallback.write(" | ".join(highlights).encode("utf-8"))
        fallback.write("\n".encode("utf-8"))
        return fallback.getvalue()
